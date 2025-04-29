import io
import base64
from typing import Dict

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import FileProcessor


class PptxProcessor(FileProcessor):
    async def process(self, content: bytes, filename: str) -> Dict:
        try:
            prs = Presentation(io.BytesIO(content))
            slides = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text = []
                tables = []
                images = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table_data = [[cell.text for cell in row.cells] for row in shape.table.rows]
                        tables.append(table_data)
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        img_bytes = image.blob
                        img_b64 = base64.b64encode(img_bytes).decode("utf-8")
                        images.append({
                            "format": image.ext,
                            "data": img_b64,
                            "width": shape.width,
                            "height": shape.height
                        })
                slides.append({
                    "page_number": slide_num,
                    "text": "\n".join(text),
                    "tables": tables,
                    "images": images
                })
            return {
                "filename": filename,
                "pages": slides,
                "page_count": len(slides)
            }
        except Exception as e:
            return {"filename": filename, "error": f"Failed to process PPTX: {str(e)}"}
